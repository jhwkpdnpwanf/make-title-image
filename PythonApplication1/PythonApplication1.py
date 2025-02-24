from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path


import fitz
import comtypes.client
import os

a = input("윗줄에 쓸 텍스트(ex 첫번째 출사이야기)")
b = input("아래에 쓸 텍스트(선유도공원)")
run_text = a + '\n' + b

def ppt_to_pdf(ppt_path, pdf_output_path):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    try:
        presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
        presentation.SaveAs(pdf_output_path, 32)  # PDF 포맷
        presentation.Close()
        print(f"PDF 저장 완료: {pdf_output_path}")
    except Exception as e:
        print(f"오류 발생: {e}")
    finally:
        powerpoint.Quit()

def crop_pdf_to_square(input_pdf, output_pdf):
    pdf_document = fitz.open(input_pdf)

    for page_number in range(len(pdf_document)):
        page = pdf_document[page_number]

        rect = page.rect
        width = rect.width
        height = rect.height
        
        square_side = height
        new_rect = fitz.Rect(0, 0, square_side, square_side)
        
        page.set_cropbox(new_rect)
    
    pdf_document.save(output_pdf)
    pdf_document.close()

    if os.path.exists(input_pdf):
        os.remove(input_pdf)

def remove_ppt(input_ppt):
    if os.path.exists(input_ppt):
        os.remove(input_ppt)

def delete_images(image_paths):
    for image_path in image_paths:
        if os.path.exists(image_path):
            os.remove(image_path)
            print(f"삭제 완료: {image_path}")


img_folder = "img_title"
img_name = "img_title1.png"  
img_path = os.path.join(img_folder, img_name) 
logo_image_path = "logo.png"

img = Image.open(img_path)

width, height = img.size

if width > height:
    left = (width - height) // 2
    upper = 0
    right = left + height
    lower = height
else:
    left = 0
    upper = (height - width) // 2
    right = width
    lower = upper + width
cropped_img = img.crop((left, upper, right, lower))

cropped_img = cropped_img.convert("RGBA")

overlay = Image.new("RGBA", cropped_img.size, (0, 0, 0, 0)) 
draw = ImageDraw.Draw(overlay)
draw.rectangle([(0, 0), cropped_img.size], fill=(0, 0, 0, 128))

combined_img = Image.alpha_composite(cropped_img, overlay)

output_path = os.path.join(img_folder, "cropped_" + img_name)
combined_img.save(output_path)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])

slide_width = prs.slide_width
slide_height = prs.slide_height

img_size = Inches(7.5) 
img_left = Inches(0) 
img_top = Inches(0) 

slide.shapes.add_picture(output_path, img_left, img_top, height=img_size)

print("검정이미지 오버레이까진 완료")

#1 흰줄
rect1_left = Inches(0.41)
rect1_top = Inches(0.41)
rect1_width = Inches(6.62)
rect1_height = Inches(6.62)

#2 흰줄
rect2_left = Inches(0.48)
rect2_top = Inches(0.48)
rect2_width = Inches(6.47)
rect2_height = Inches(6.47)

num_steps = 10
left_step = (rect2_left - rect1_left) / num_steps
top_step = (rect2_top - rect1_top) / num_steps
width_step = (rect2_width - rect1_width) / num_steps
height_step = (rect2_height - rect1_height) / num_steps

for i in range(num_steps + 1):
    left = rect1_left + left_step * i
    top = rect1_top + top_step * i
    width = rect1_width + width_step * i
    height = rect1_height + height_step * i

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=left,
        top=top,
        width=width,
        height=height
    )
    shape.fill.background()  # 내부 투명
    shape.line.color.rgb = RGBColor(255, 255, 255)  # 흰색 테두리
    shape.line.width = Pt(1)  # 테두리 두께
    shape.shadow.inherit = False
print("테두리 박스 생성 중")

# 글씨쓰기
textbox_width = Inches(7.5) 
textbox_height = Inches(2.0) 
left = 0 
top = Inches(3)

textbox = slide.shapes.add_textbox(left, top, textbox_width, textbox_height)

text_frame = textbox.text_frame

# 텍스트 설정
p = text_frame.paragraphs[0]
p.alignment = PP_ALIGN.CENTER 
run = p.add_run()
run.text = run_text
font = run.font
font.name = "맑은 고딕"  # 글꼴
font.size = Pt(42)       # 크기
font.bold = True         # 볼드체
font.color.rgb = RGBColor(255, 255, 255)

textbox.fill.background()

output_ppt = "photo.pptx"
prs.save(output_ppt)


# 로고
left = Inches(3.10)
top = Inches(0.85)
width = Inches(1.30)
height = Inches(1.30) 

slide.shapes.add_picture(logo_image_path, left, top, width=width, height=height)

print(f"로고가 추가완료")

print("최종 PPT 완성")
print("이미지 생성중")

output_ppt = "photo.pptx"
prs.save(output_ppt)


ppt_path = os.path.abspath("photo.pptx")
pdf_output_path = os.path.abspath("photo.pdf") 

if not os.path.exists(ppt_path):
    raise FileNotFoundError(f"PowerPoint 파일을 찾을 수 없습니다: {ppt_path}")

ppt_to_pdf(ppt_path, pdf_output_path)
print("PDF로 변환완료")

input_pdf = "photo.pdf"
output_pdf = "img_title0.pdf"


crop_pdf_to_square(input_pdf, output_pdf)
print("PDF로 변환완료")


pdf_document = fitz.open(output_pdf)

page = pdf_document[0]
pix = page.get_pixmap()
pix.save("title.png")

print("PDF 첫 번째 페이지를 'title.png'로 저장했습니다.")
pdf_document.close()


input_ppt = "photo.pptx"
remove_ppt(input_ppt)


print("작업완료")