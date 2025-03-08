import os
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw, ImageFont
from pdf2image import convert_from_path
import fitz
import comtypes.client

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

a = input("윗줄에 쓸 텍스트(ex 첫번째 출사이야기)")
b = input("아래에 쓸 텍스트(선유도공원)")
run_text = a + '\n' + b

img_folder = os.path.join(BASE_DIR, "img_title")
logo_image_path = os.path.join(BASE_DIR, "assets", "logo.png")

output_folder = os.path.join(BASE_DIR, "output")
os.makedirs(output_folder, exist_ok=True)

output_ppt = os.path.join(output_folder, "photo.pptx")
pdf_output_path = os.path.join(output_folder, "photo.pdf")

def ppt_to_pdf(ppt_path, pdf_output_path):
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1
    try:
        presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
        presentation.SaveAs(pdf_output_path, 32)
        presentation.Close()
    except Exception as e:
        print(f"오류 발생: {e}")
    finally:
        powerpoint.Quit()

def crop_pdf_to_square(input_pdf, output_pdf):
    pdf_document = fitz.open(input_pdf)
    for page_number in range(len(pdf_document)):
        page = pdf_document[page_number]
        rect = page.rect
        width, height = rect.width, rect.height
        square_side = height
        new_rect = fitz.Rect(0, 0, square_side, square_side)
        page.set_cropbox(new_rect)
    pdf_document.save(output_pdf)
    pdf_document.close()

img_name = next((file for file in os.listdir(img_folder) if file.lower().endswith(('.png', '.jpg', '.jpeg'))), None)

if img_name:
    img_path = os.path.join(img_folder, img_name)

    img = Image.open(img_path)
    min_side = min(img.size)
    left = (img.width - min_side) // 2
    top = (img.height - min_side) // 2
    right = left + min_side
    bottom = top + min_side

    cropped_img = img.crop((left, top, right, bottom))

    # 검정 오버레이
    cropped_img = cropped_img.convert("RGBA")
    overlay = Image.new("RGBA", cropped_img.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    draw.rectangle([(0, 0), cropped_img.size], fill=(0, 0, 0, 128))

    combined_img = Image.alpha_composite(cropped_img, overlay)

    output_path = os.path.join(img_folder, "cropped_" + img_name)
    combined_img.save(output_path)

    prs = Presentation()

    # 슬라이드 정사각형
    prs.slide_width = Inches(7.5)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.add_picture(output_path, Inches(0), Inches(0), height=Inches(7.5))

    # 흰색 테두리
    rect1_left = Inches(0.41)
    rect1_top = Inches(0.41)
    rect1_width = Inches(6.62)
    rect1_height = Inches(6.62)

    rect2_left = Inches(0.48)
    rect2_top = Inches(0.48)
    rect2_width = Inches(6.47)
    rect2_height = Inches(6.47)

    num_steps = 10
    left_step = (rect2_left - rect1_left) / num_steps
    top_step = (rect2_top - rect1_top) / num_steps
    width_step = (rect2_width - rect1_width) / num_steps
    height_step = (rect2_height - rect1_height) / num_steps

    for i in range(num_steps + 1):
        left = rect1_left + left_step * i
        top = rect1_top + top_step * i
        width = rect1_width + width_step * i
        height = rect1_height + height_step * i

        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left=left,
            top=top,
            width=width,
            height=height
        )
        shape.fill.background()
        shape.line.color.rgb = RGBColor(255, 255, 255)
        shape.line.width = Pt(1)
        shape.shadow.inherit = False

    textbox_width = Inches(7.5)
    textbox_height = Inches(2.0)
    left = 0
    top = Inches(3)

    textbox = slide.shapes.add_textbox(left, top, textbox_width, textbox_height)
    text_frame = textbox.text_frame

    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = run_text
    font = run.font
    font.name = "맑은 고딕"
    font.size = Pt(42)
    font.bold = True
    font.color.rgb = RGBColor(255, 255, 255)

    textbox.fill.background()

    try:
        slide.shapes.add_picture(logo_image_path, Inches(3.10), Inches(0.85), width=Inches(1.30), height=Inches(1.30))
        print("로고됨")
    except:
        print("로고안됨")

    prs.save(output_ppt)
    ppt_to_pdf(output_ppt, pdf_output_path)

    if os.path.exists(output_path):
        os.remove(output_path)

    if os.path.exists(output_ppt):
        os.remove(output_ppt)

    print("작업완료")
else:
    print("img_title 폴더에 이미지 파일이 없습니다.")
